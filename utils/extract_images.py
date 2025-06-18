from pptx import Presentation
import argparse
import os

def extract_images_from_shape(shape, output_dir, slide_index, img_index):
    if shape.shape_type == 13:  
        image = shape.image
        ext = image.ext
        img_bytes = image.blob
        filename = f"slide_{slide_index}_img_{img_index}.{ext}"
        with open(os.path.join(output_dir, filename), "wb") as f:
            f.write(img_bytes)
        return 1
    elif shape.shape_type == 6: 
        count = 0
        for subshape in shape.shapes:
            count += extract_images_from_shape(subshape, output_dir, slide_index, img_index + count)
        return count
    return 0

def extract_images_from_slide(path, dest):
    prs = Presentation(path)
    img_count = 0
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            img_count += extract_images_from_shape(shape, dest, i + 1, img_count + 1)
    print(f"✅ Đã trích xuất {img_count} ảnh vào thư mục {dest}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract images from PowerPoint slides.")
    parser.add_argument("pptx_path", help="Path to the PowerPoint file")
    parser.add_argument("output_dir", help="Directory to save extracted images")
    args = parser.parse_args()
    if not os.path.exists(args.output_dir):
        os.makedirs(args.output_dir)
    extract_images_from_slide(args.pptx_path, args.output_dir)


