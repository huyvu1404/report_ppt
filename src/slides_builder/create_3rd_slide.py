from .slide_utils import *
from slides_builder import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from charts_generator import *
from utils import *
from llm import *
       
def create_third_slide(prs, current_data, main_topic, current_json_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    left, top, width, height = Inches(2.31), Inches(0.13), Inches(1.42), Inches(0.2)
    spacing = Inches(0.16)
    texts = ["TÓM TẮT", "SO SÁNH THẢO LUẬN", "TÓM TẮT THẢO LUẬN"]
    colors = [RGBColor(220, 220, 220), RGBColor(255, 243, 205), RGBColor(220, 220, 220)]
    text_colors = [RGBColor(198, 198, 198), RGBColor(156, 91, 205), RGBColor(198, 198, 198)]
    width_scale = [1, 1.1, 1]
    create_section_rectangle(shapes, 3, (left, top, width, height), colors, text_colors, width_scale, spacing=spacing, texts=texts)

    left, top, width, height = Inches(0), Inches(0.46), Inches(10), Inches(0.39)
    colors = [RGBColor(102, 14, 207), RGBColor(29, 191, 142)]
    create_rectangle(shapes, (left, top, width, height), gradient=True, colors=colors, shadow=False)
    
    left, top, width, height = Inches(-0.35), Inches(0.45), Inches(0.77), Inches(0.4)
    create_rounded_rectangle( shapes, (left, top, width, height), adjustment=0.5, color=RGBColor(255, 255, 255), shadow=False)
    
    left, top, width, height = Inches(0.52), Inches(0.49), Inches(5.90), Inches(0.34)
    texts = "NỘI DUNG THẢO LUẬN NỔI BẬT VỀ SHB VÀ ĐỐI THỦ"
    create_text_box(shapes, texts, (left, top, width, height), fontsize=Pt(14), fontweight="bold", text_color=RGBColor(255, 255, 255), text_alignment=1)

    left, top, width, height = Inches(0.14), Inches(0.89), Inches(9.82), Inches(1.2)
#     texts = """Hoạt động fanpage tạo thảo luận tích cực nổi bật tại các ngân hàng Techcombank và VPBank nhờ hoạt động tương tác xúc tiến chiến dịch như minigame và livestream. Cùng với đó, choỗi hoạt động săn vé sự kiện cũng ghi nhận bàn luận sôi nổi.
# Thảo luận tiêu cực của tuần vẫn xoay quanh lỗi hệ thống săn vé sự kiện trên ứng dụng và website chiến dịch, tiêu biểu đến từ VPBank và Techcombank."""
    transformed_json = prepare_json_data_3rd(current_json_data)
    insight = get_third_insight(transformed_json)
    create_rounded_rectangle(shapes, (left, top, width, height), adjustment=0.1, color=RGBColor(255, 243, 205), texts=insight, fontsize=Pt(9), text_color=RGBColor(0, 0, 0), text_alignment=1, shadow=True)

    data = prepare_table_data(current_data, main_topic=main_topic, rows='Labels1',columns=['Topic', 'Sentiment'], values='Id', aggfunc='count')    
    chart, topics, labels, width_ratios = generate_table(data)
    width_ratios[0] = 1
    width = Inches(9.82) * len(topics) / 8
    bar_left = Inches(0.22) + (Inches(9.82) - width) / 2
    shapes.add_picture(chart, bar_left, Inches(2.21), width, Inches(3.42))
    # shapes.add_picture('./img/logo-slide3.png', Inches(1.6), Inches(2.1), width=Inches(7.91), height=Inches(0.27))

    top = Inches(2.18)
    current_left = Inches(1.66) + bar_left
    for i, topic in enumerate(topics):
        topic = topic.replace(" ", "")
        
        width, height = LOGO_SIZES.get(topic)
        left = current_left - width / 2
       
        shapes.add_picture(f'src/assets/icons/{topic}.png', left, top, width, height)
        if i < len(width_ratios) - 1:
            current_left += Inches((width_ratios[i] + width_ratios[i+1]) * 1.05 / 2)

    left, top, width, height = bar_left - Inches(0.06), Inches(2.39), Inches(1.18), Inches(0.2)
    distance = Inches(0.08)
    for i in range(11):
        create_rectangle(shapes, (left, top + i * (height + distance), width, height),texts=labels[i], text_color=RGBColor(0,0,0), fontsize=Pt(5.5), color=RGBColor(254, 219, 106), shadow=False, fontweight="bold", text_alignment=2)
