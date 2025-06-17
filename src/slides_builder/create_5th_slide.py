from .slide_utils import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from charts_generator import prepare_stacked_bar_data_5th, generate_stacked_bar_chart_5th
from llm import prepare_json_data_5th, get_fifth_insight
       
def create_fifth_slide(prs, current_data, main_topic, current_json_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    left, top, width, height = Inches(2.31), Inches(0.13), Inches(1.42), Inches(0.2)
    spacing = Inches(0.16)
    texts = ["TÓM TẮT", "SO SÁNH THẢO LUẬN", "TÓM TẮT THẢO LUẬN"]
    colors = [RGBColor(220, 220, 220), RGBColor(220, 220, 220),  RGBColor(255, 243, 205)]
    text_colors = [RGBColor(198, 198, 198), RGBColor(198, 198, 198), RGBColor(156, 91, 205)]
    width_scale = [1, 1, 1.1]
    create_section_rectangle(shapes, 3, (left, top, width, height), colors, text_colors, width_scale, spacing=spacing, texts=texts)

    left, top, width, height = Inches(0), Inches(0.46), Inches(10), Inches(0.39)
    colors = [RGBColor(102, 14, 207), RGBColor(29, 191, 142)]
    create_rectangle(shapes, (left, top, width, height), gradient=True, colors=colors, shadow=False)
    
    left, top, width, height = Inches(-0.35), Inches(0.45), Inches(0.77), Inches(0.4)
    create_rounded_rectangle( shapes, (left, top, width, height), adjustment=0.5, color=RGBColor(255, 255, 255), shadow=False)
    
    left, top, width, height = Inches(0.52), Inches(0.49), Inches(5.01), Inches(0.34)
    texts = "TOP CÁC NỘI DUNG THẢO LUẬN VỀ SHB"
    create_text_box(shapes, texts, (left, top, width, height), fontsize=Pt(14), fontweight="bold", text_color=RGBColor(255, 255, 255), text_alignment=1)

    left, top, width, height = Inches(0), Inches(0.92), Inches(10), Inches(1.77)
    transformed_json = prepare_json_data_5th(current_json_data, main_topic)
    insight = get_fifth_insight(transformed_json)
    create_text_box(shapes, insight, (left, top, width, height), fontsize=Pt(9), text_color=RGBColor(0, 0, 0), text_alignment=1, shadow=True, adjustment=0.05, color=RGBColor(255, 243, 205))
    

    data = prepare_stacked_bar_data_5th(current_data, main_topic=main_topic, 
        rows='Labels1',
        columns=['Topic', 'Sentiment'], 
        values='Id', 
        aggfunc='count'
        )
    chart = generate_stacked_bar_chart_5th(data)
    shapes.add_picture(chart, Inches(0), Inches(2.90), Inches(9.74), Inches(2.64))    