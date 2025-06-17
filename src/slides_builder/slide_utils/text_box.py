from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_UNDERLINE

import re


def create_text_box(shapes, texts, attribute, **kwargs):
    left, top, width, height = attribute
    textbox = shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.alignment = kwargs.get('text_alignment', PP_ALIGN.LEFT)

    font_size = kwargs.get('fontsize', Pt(8))
    font_bold = kwargs.get('fontweight') == "bold"
    font_italic = kwargs.get('fontstyle') == "italic"
    font_color = kwargs.get('text_color', RGBColor(0, 0, 0))

    # Xử lý văn bản chứa **...**
    parts = re.split(r'(\*\*.*?\*\*)', texts)

    for part in parts:
        run = p.add_run()
        if part.startswith("**") and part.endswith("**"):
            content = part[2:-2]
            run.text = content.upper()  # Nếu không muốn in hoa thì dùng: content
            run.font.bold = True
        else:
            run.text = part
            run.font.bold = font_bold

        run.font.size = font_size
        run.font.italic = font_italic
        run.font.color.rgb = font_color

def create_text_box_with_url(shapes, texts, attribute, url_label=None, url=None, label_color=RGBColor(254, 219, 106),**kwargs):
    left, top, width, height = attribute
    textbox = shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    
    # Run chính
    run_main = p.add_run()
    run_main.text = texts
    run_main.font.size = kwargs.get('fontsize', Pt(12))
    run_main.font.color.rgb = kwargs.get('text_color', RGBColor(0, 0, 0))
    if kwargs.get('fontweight') == "bold":
        run_main.font.bold = True
    if kwargs.get('fontstyle') == "italic":
        run_main.font.italic = True
    if url and url_label:
        run_link = p.add_run()
        run_link.text = f" {url_label}"  
        run_link.hyperlink.address = url
        run_link.font.size = kwargs.get('fontsize', Pt(12))
        run_link.font.color.rgb = label_color
        run_link.font.underline = MSO_UNDERLINE.DASH_LINE
        if kwargs.get('fontweight') == "bold":
            run_link.font.bold = True

    p.alignment = kwargs.get('text_alignment', PP_ALIGN.LEFT)

