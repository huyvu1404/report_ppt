from .slide_utils import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from utils import LOGO_SIZES, PROJECT_DIR
from llm import prepare_json_data_2nd, get_second_insight
from charts_generator import prepare_stacked_bar_data_2nd, generate_stacked_bar_chart_2nd

def create_second_slide(prs, current_data, previous_data, main_topic, current_json, previous_json):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    left, top, width, height = Inches(2.31), Inches(0.13), Inches(1.42), Inches(0.2)
    spacing = Inches(0.16)
    texts = ["TÓM TẮT", "SO SÁNH THẢO LUẬN", "TÓM TẮT THẢO LUẬN"]
    colors = [RGBColor(220, 220, 220), RGBColor(255, 243, 205), RGBColor(220, 220, 220)]
    text_colors = [RGBColor(198, 198, 198), RGBColor(156, 91, 205), RGBColor(198, 198, 198)]
    width_scale = [1, 1.1, 1]
    create_section_rectangle(shapes, 3, (left, top, width, height), colors, text_colors, width_scale, spacing=spacing, texts=texts)

    left, top, width, height = Inches(0), Inches(0.46), Inches(10), Inches(0.39)
    colors = [RGBColor(102, 14, 207), RGBColor(29, 191, 142)]
    create_rectangle(shapes, (left, top, width, height), gradient=True, colors=colors, shadow=False)
    
    left, top, width, height = Inches(-0.35), Inches(0.45), Inches(0.77), Inches(0.4)
    create_rounded_rectangle( shapes, (left, top, width, height), adjustment=0.5, color=RGBColor(255, 255, 255), shadow=False)
    
    left, top, width, height = Inches(0.52), Inches(0.49), Inches(5.01), Inches(0.34)
    texts = "SẮC THÁI THẢO LUẬN CỦA SHB VÀ ĐỐI THỦ"
    create_text_box(shapes, texts, (left, top, width, height), fontsize=Pt(14), fontweight="bold", text_color=RGBColor(255, 255, 255), text_alignment=1)

    left, top, width, height = Inches(0.14), Inches(0.96), Inches(9.82), Inches(1.39)
    
    transformed_json = prepare_json_data_2nd(current_json, previous_json)
    insight = get_second_insight(transformed_json)
    create_rounded_rectangle(shapes, (left, top, width, height), adjustment=0.05, color=RGBColor(255, 243, 205), texts=insight, fontsize=Pt(9), text_color=RGBColor(0, 0, 0), text_alignment=1, shadow=True)
    
    left, top, width, height = Inches(0.14), Inches(2.52), Inches(9.74), Inches(2.86)
    create_rounded_rectangle(shapes, (left, top, width, height), adjustment=0.05, color=RGBColor(255, 255, 255), shadow=True, line_color=RGBColor(227,227,227))
    
    data = prepare_stacked_bar_data_2nd(current_data, previous_data, main_topic=main_topic, rows='Sentiment', columns='Topic', values='Id', aggfunc='count')
    chart, topics = generate_stacked_bar_chart_2nd(data)
    width = Inches(9.59) * len(topics) / 8
    bar_left = Inches(0.14) + (Inches(9.59) - width) / 2
    shapes.add_picture(chart, bar_left, Inches(2.5), width, Inches(2.5))
    
    spacing = Inches(1.17)
    top = Inches(5.01)
    based_left = bar_left + Inches(0.66)
    for i, topic in enumerate(topics):
        topic = topic.replace(" ", "")
        width, height = LOGO_SIZES.get(topic)
        left = based_left + i * spacing - width / 2
        ICON_PATH = PROJECT_DIR / "assets/icons" / f"{topic}.png"
        shapes.add_picture(str(ICON_PATH), left, top, width, height)
