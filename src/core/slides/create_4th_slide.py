from .slide_utils import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from core.charts import prepare_nested_data, generate_nested_chart
from core.insights import prepare_json_data_4th, get_fourth_insight
       
def create_fourth_slide(prs, current_data, previous_data, main_topic, current_json_data, previous_json_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    left, top, width, height = Inches(2.31), Inches(0.13), Inches(1.42), Inches(0.2)
    spacing = Inches(0.16)
    texts = ["TÓM TẮT", "SO SÁNH THẢO LUẬN", "TÓM TẮT THẢO LUẬN"]
    colors = [RGBColor(220, 220, 220), RGBColor(220, 220, 220),  RGBColor(255, 243, 205)]
    text_colors = [RGBColor(198, 198, 198), RGBColor(198, 198, 198), RGBColor(156, 91, 205)]
    width_scale = [1, 1, 1.1]
    create_section_rectangle(shapes, 3, (left, top, width, height), colors, text_colors, width_scale, spacing=spacing, texts=texts)

    left, top, width, height = Inches(0), Inches(0.46), Inches(10), Inches(0.39)
    colors = [RGBColor(102, 14, 207), RGBColor(29, 191, 142)]
    create_rectangle(shapes, (left, top, width, height), gradient=True, colors=colors, shadow=False)
    
    left, top, width, height = Inches(-0.35), Inches(0.45), Inches(0.77), Inches(0.4)
    create_rounded_rectangle( shapes, (left, top, width, height), adjustment=0.5, color=RGBColor(255, 255, 255), shadow=False)
    
    left, top, width, height = Inches(0.52), Inches(0.49), Inches(5.01), Inches(0.34)
    texts = "TỔNG QUAN THẢO LUẬN VỀ SHB"
    create_text_box(shapes, texts, (left, top, width, height), fontsize=Pt(14), fontweight="bold", text_color=RGBColor(255, 255, 255), text_alignment=1)

    left, top, width, height = Inches(0.00), Inches(1.11), Inches(10.00), Inches(1.03)

    transformed_json = prepare_json_data_4th(current_json_data, previous_json_data, main_topic)
    insight = get_fourth_insight(transformed_json) 
    create_text_box(shapes, insight, (left, top, width, height), fontsize=Pt(9), text_color=RGBColor(0, 0, 0), text_alignment=1, shadow=True, adjustment=0.05, color=RGBColor(255, 243, 205))
    
    create_rectangle(shapes, (Inches(0), Inches(2.48), Inches(10), Inches(2.97)), color=RGBColor(255, 255, 255), shadow=True)
    
    current_sentiment_data, current_channel_data, previous_sentiment_data, previous_channel_data = prepare_nested_data(current_data, previous_data, main_topic=main_topic)
    nested_bar_chart = generate_nested_chart(current_sentiment_data, current_channel_data, previous_sentiment_data, previous_channel_data)
    if nested_bar_chart:
        shapes.add_picture(nested_bar_chart, Inches(0.3), Inches(2.54), width=Inches(9.21), height=Inches(2.86))
        texts = "TỈ TRỌNG THẢO LUẬN TRÊN CÁC KÊNH TRỰC TUYẾN VÀ SẮC THÁI THẢO LUẬN CỦA SHB"
        create_rectangle(shapes, (Inches(3.17), Inches(2.46), Inches(3.47), Inches(0.40)), color=RGBColor(246, 146, 0), shadow=False, texts=texts, fontsize=Pt(9), text_color=RGBColor(255, 255, 255), text_alignment=2, fontweight="bold")
        create_text_box(shapes, "Tuần này", (Inches(1.46), Inches(4.92), Inches(2.04), Inches(0.44)), fontsize=Pt(10), text_color=RGBColor(0, 0, 0), text_alignment=2, fontstyle="italic", fontweight="bold")
        create_text_box(shapes, "Tuần trước", (Inches(6.26), Inches(4.92), Inches(2.04), Inches(0.44)), fontsize=Pt(10), text_color=RGBColor(196, 189, 151), text_alignment=2, fontstyle="italic")