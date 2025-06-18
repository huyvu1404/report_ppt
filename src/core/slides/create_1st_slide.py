from .slide_utils import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from core.charts import prepare_doughnut_data, generate_doughnut_chart, prepare_stacked_bar_data, generate_stacked_bar_chart
from constants import LOGO_SIZES, PROJECT_DIR
from core.insights import prepare_json_data, get_first_insight
from .create_slide_header import create_slide_header

       
def create_first_slide(prs, current_data, previous_data, main_topic, current_json, previous_json):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    texts = ["TÓM TẮT", "SO SÁNH THẢO LUẬN", "TÓM TẮT THẢO LUẬN"]
    shape_colors = [RGBColor(220, 220, 220), RGBColor(255, 243, 205), RGBColor(220, 220, 220)]
    text_colors = [RGBColor(198, 198, 198), RGBColor(156, 91, 205), RGBColor(198, 198, 198)]
    shape_width_scale = [1, 1.1, 1]
    create_slide_header(shapes, texts, shape_colors, text_colors, shape_width_scale)
    
    title_left, title_top, title_width, title_height = Inches(0.52), Inches(0.49), Inches(5.90), Inches(0.34)
    texts = "TỔNG QUAN THẢO LUẬN VỀ SHB VÀ ĐỐI THỦ"
    create_text_box(shapes, texts, (title_left, title_top, title_width, title_height), fontsize=Pt(14), fontweight="bold", text_color=RGBColor(255, 255, 255), text_alignment=1)

    insight_left, insight_top, insight_width, insight_height = Inches(0.07), Inches(0.9), Inches(4.61), Inches(2.52)
    transformed_json_data = prepare_json_data(current_json, previous_json)
  
    insight = get_first_insight(transformed_json_data) 
    create_rounded_rectangle(shapes, (insight_left, insight_top, insight_width, insight_height), adjustment=0.05, color=RGBColor(255, 243, 205), texts=insight, fontsize=Pt(9), text_color=RGBColor(0, 0, 0), text_alignment=1, shadow=True)
    
    doughnut_left, doughnut_top, doughnut_width, doughnut_height = Inches(4.74), Inches(1.08), Inches(4.67), Inches(2.52)
    
    current_data_pie, previous_data_pie = prepare_doughnut_data(current_data, previous_data, main_topic="SHB",columns='Topic', values='Id', aggfunc='count' )
    doughnut_chart = generate_doughnut_chart(current_data_pie, previous_data_pie) 
    shapes.add_picture(doughnut_chart, doughnut_left, doughnut_top, doughnut_width, doughnut_height)
    
    left, top, width, height = Inches(5.83), Inches(0.97), Inches(3), Inches(0.24)
    title = "THỊ PHẦN THẢO LUẬN CỦA SHB VÀ ĐỐI THỦ"
    create_rectangle(shapes, (left, top, width, height), color=RGBColor(0, 112, 192), shadow=False, fontweight="bold", texts=title, fontsize=Pt(8), text_color=RGBColor(255, 255, 255), text_alignment=2)
    
    left, top, width, height = Inches(0.14), Inches(3.49), Inches(9.7), Inches(2.08)
    create_rounded_rectangle(shapes, (left, top, width, height), adjustment=0.05, color=RGBColor(255, 255, 255), shadow=True, line_color=RGBColor(227,227,227))
    
    current_data_bar, previous_data_bar = prepare_stacked_bar_data(current_data, previous_data, main_topic=main_topic, rows='Channel', columns='Topic', values='Id', aggfunc='count')
    stacked_bar, topics = generate_stacked_bar_chart(current_data_bar, previous_data_bar)
    width = Inches(9.24) * len(topics) / 8
    stacked_bar_left = Inches(0.29) + (Inches(9.24) - width) / 2
    if stacked_bar:
        shapes.add_picture(stacked_bar, stacked_bar_left, Inches(3.57), width=width, height=Inches(1.77))
        spacing = Inches(1.17)
        top = Inches(5.3)
        based_left = stacked_bar_left + Inches(0.51)
        for i, topic in enumerate(topics):
            topic = topic.replace(" ", "")
            width, height = LOGO_SIZES.get(topic)
            left = based_left + i * spacing - width / 2
            ICON_PATH = PROJECT_DIR / "assets/icons" / f"{topic}.png"

            shapes.add_picture(str(ICON_PATH), left, top, width, height)

        left, top, width, height = Inches(0.41), Inches(3.55), Inches(3.5), Inches(0.22)
        title = "KÊNH THẢO LUẬN VỀ SHB VÀ CÁC NGÂN HÀNG ĐỐI THỦ"
        create_rectangle(shapes, (left, top, width, height), color=RGBColor(246, 146, 0), shadow=False, fontweight="bold", texts=title, fontsize=Pt(7), text_color=RGBColor(255, 255, 255), text_alignment=2)
        

